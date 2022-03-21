from morphy_using import name_change
from pptx import Presentation

def PPTX_GENERATOR(data):
    try:
        data['name'] = name_change(data['name'], data['case'])
    except KeyError:
        pass
    prs = Presentation(data['template'])
    for shape in prs.slides[0].shapes:  # Перебираем объекты на слайде
        if shape.has_text_frame:  # Если у объекта есть текст,
            for key, value in data.items():  # То смотрим равен является ли он словом-заместителем
                try:
                    # Меняем слово-заместитель в презентации на значение этого слова в нашем словаре data
                    if key.lower() == shape.text_frame.text.lower():
                        shape.text_frame.paragraphs[0].runs[0].text = str(value)
                except AttributeError:
                    pass

    # указываем создателя
    prs.core_properties.author = ""
    # заголовок файла (ВАЖНО ДЛЯ PDF он в заголовке пишется)
    prs.core_properties.title = ""

    #  для хранения сертификатов создаём две папки, в которых тоже будут папки
    #  внитри GENERATED_PPTX и GENERATED_PDF будут папки-даты
    prs.save('GENERATED_PPTX/' + data['date'] + '/' + data['file_name'] + '_' + data['id'] + '.pptx')
    # имя файла для перевода его в PDF
    return (data['file_name'] + '_' + data['id'])

# # Дублируем слайд
# def duplicate_slide(pres, index=0):
#     template = pres.slides[index]
#     blank_slide_layout = pres.slide_layouts[0] # Говорим что слайд будет титульным (1 - титульный слайд)

#     copied_slide = pres.slides.add_slide(blank_slide_layout)
    

#     for shp in template.shapes:
#         el = shp.element
#         newel = copy.deepcopy(el)
#         copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

#     # Копируем .rels которые ведут нас к изображению, без этого картинка не отоброзится
#     for _, value in six.iteritems(template.part.rels):
#         # Make sure we don't copy a notesSlide relation as that won't exist
#         if "notesSlide" not in value.reltype:
#             copied_slide.part.rels.add_relationship(
#                 value.reltype,
#                 value._target,
#                 value.rId
#             )